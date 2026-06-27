import logging
from pathlib import Path

from app.core.config import settings

logger = logging.getLogger(__name__)


class BaseParser:
    name = "base_parser"

    def parse(self, storage_path: str, fallback_title: str) -> str:
        path = Path(storage_path)
        if path.exists():
            return path.read_text(encoding="utf-8", errors="replace")
        return f"{fallback_title}\nSource path: {storage_path}"


class PDFParser(BaseParser):
    name = "pdf_parser"

    def parse(self, storage_path: str, fallback_title: str) -> str:
        path = Path(storage_path)
        if not path.exists():
            return f"{fallback_title}\nFile not found: {storage_path}"

        texts: list[str] = []
        try:
            import pymupdf
            doc = pymupdf.open(str(path))
            for page in doc:
                page_text = page.get_text()
                if page_text.strip():
                    texts.append(page_text)
            doc.close()
            if texts:
                return "\n\n".join(texts)
        except Exception as e:
            logger.warning("pymupdf failed for %s: %s", storage_path, e)

        try:
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        texts.append(page_text)
            if texts:
                return "\n\n".join(texts)
        except Exception as e:
            logger.warning("pdfplumber failed for %s: %s", storage_path, e)

        return f"{fallback_title}\nPDF parsing failed. Path: {storage_path}"


class WordParser(BaseParser):
    name = "word_parser"

    def parse(self, storage_path: str, fallback_title: str) -> str:
        path = Path(storage_path)
        if not path.exists():
            return f"{fallback_title}\nFile not found: {storage_path}"

        try:
            from docx import Document
            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs) if paragraphs else fallback_title
        except Exception as e:
            logger.warning("python-docx failed for %s: %s", storage_path, e)
            return f"{fallback_title}\nWord parsing failed. Path: {storage_path}"


class ExcelParser(BaseParser):
    name = "excel_parser"

    def parse(self, storage_path: str, fallback_title: str) -> str:
        path = Path(storage_path)
        if not path.exists():
            return f"{fallback_title}\nFile not found: {storage_path}"

        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            rows: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows.append(f"--- Sheet: {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append("\t".join(cells))
            wb.close()
            return "\n".join(rows) if rows else fallback_title
        except Exception as e:
            logger.warning("openpyxl failed for %s: %s", storage_path, e)
            return f"{fallback_title}\nExcel parsing failed. Path: {storage_path}"


class PowerPointParser(BaseParser):
    name = "powerpoint_parser"

    def parse(self, storage_path: str, fallback_title: str) -> str:
        path = Path(storage_path)
        if not path.exists():
            return f"{fallback_title}\nFile not found: {storage_path}"

        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slides: list[str] = []
            for idx, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            if p.text.strip():
                                slide_texts.append(p.text)
                if slide_texts:
                    slides.append(f"Slide {idx}: " + " ".join(slide_texts))
            return "\n\n".join(slides) if slides else fallback_title
        except Exception as e:
            logger.warning("python-pptx failed for %s: %s", storage_path, e)
            return f"{fallback_title}\nPowerPoint parsing failed. Path: {storage_path}"


class MarkdownParser(BaseParser):
    name = "markdown_parser"


class TextParser(BaseParser):
    name = "text_parser"


class CSVParser(BaseParser):
    name = "csv_parser"


class DocumentParser:
    text_extensions = {"txt", "md", "csv"}
    office_extensions = {"docx", "pptx", "xlsx"}
    pdf_extensions = {"pdf"}
    strategies = {
        "pdf": PDFParser(),
        "docx": WordParser(),
        "xlsx": ExcelParser(),
        "pptx": PowerPointParser(),
        "md": MarkdownParser(),
        "txt": TextParser(),
        "csv": CSVParser(),
    }

    def parse(self, *, storage_path: str, file_type: str | None, fallback_title: str) -> dict:
        extension = (file_type or Path(storage_path).suffix.lstrip(".")).lower()
        parser = self.strategies.get(extension, BaseParser())
        text = parser.parse(storage_path, fallback_title)

        return {
            "text": text,
            "pages": [{"page": 1, "text": text}],
            "parser": parser.name,
        }
